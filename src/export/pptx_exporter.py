"""PowerPoint エクスポート: 6枚構成の経営報告スライド"""
import io
from datetime import date
from typing import List, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.export._helpers import (
    INGREDIENTS, load_sku_fc_df, next_period, get_actual_sellout_12m,
    BRAND_BLUE, BRAND_RED, BRAND_GREEN, BRAND_GRAY,
)
from src.forecast.market_forecast import aggregate_market, MarketForecaster
from src.forecast.penetration import PenetrationForecaster
from src.forecast.scenario import build_penetration_events
from src.forecast.currency import convert_amounts

plt.rcParams["font.family"] = ["Hiragino Sans", "AppleGothic", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

_SLIDE_W = Inches(10)
_SLIDE_H = Inches(7.5)

_C_BLUE  = RGBColor(0x25, 0x63, 0xEB)
_C_RED   = RGBColor(0xDC, 0x26, 0x26)
_C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_C_DARK  = RGBColor(0x1E, 0x29, 0x3B)
_C_GRAY  = RGBColor(0x94, 0xA3, 0xB8)


def _rgb_hex(hex_str: str) -> tuple:
    h = hex_str.lstrip("#")
    return tuple(int(h[i:i+2], 16) / 255 for i in (0, 2, 4))


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title_bar(slide, text: str, subtitle: str = ""):
    tf = slide.shapes.add_textbox(Inches(0), Inches(0), _SLIDE_W, Inches(0.8))
    tf.fill.solid()
    tf.fill.fore_color.rgb = _C_BLUE
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold  = True
    run.font.size  = Pt(20)
    run.font.color.rgb = _C_WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.text_frame.margin_left = Inches(0.2)
    tf.text_frame.margin_top  = Inches(0.12)

    if subtitle:
        tf2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.82), Inches(8), Inches(0.3))
        tf2.text_frame.paragraphs[0].text = subtitle
        tf2.text_frame.paragraphs[0].runs[0].font.size  = Pt(9)
        tf2.text_frame.paragraphs[0].runs[0].font.color.rgb = _C_GRAY


def _add_textbox(slide, text: str, x, y, w, h, size=11, bold=False, color=None):
    tf = slide.shapes.add_textbox(x, y, w, h)
    for para_text in text.split("\n"):
        p = tf.text_frame.add_paragraph() if tf.text_frame.paragraphs[0].text else tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if para_text != text.split("\n")[0]:
            p = tf.text_frame.add_paragraph()
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(size)
            run.font.bold = bold
    return tf


def _fig_to_bytes(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


# ── スライド生成関数 ──────────────────────────────────────────────

def _slide_title(prs, ing_name, axis2, currency, today):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_textbox(Inches(0), Inches(0), _SLIDE_W, _SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _C_BLUE

    tf = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.6))
    p = tf.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "GE/BS 販売予測レポート"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = _C_WHITE

    sub = tf.text_frame.add_paragraph()
    r2 = sub.add_run()
    r2.text = f"{ing_name}  |  {today}  |  通貨: {currency}  |  シナリオ幅: {axis2}"
    r2.font.size = Pt(14)
    r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)


def _slide_summary(prs, df_adj, df_exc, currency, fx, ing_name, axis2):
    slide = _blank_slide(prs)
    unit = "万JPY" if currency == "JPY" else "万USD"
    _add_title_bar(slide, "① サマリー", f"{ing_name} — 36ヶ月予測サマリー（制度変更加味・{axis2}）")

    df = df_adj.copy()
    df["year"]   = df["period"].str[:4]
    df["amount"] = convert_amounts(df["amount_jpy"].values, currency, fx) / 10_000
    yr_totals    = df.groupby("year")["amount"].sum()
    years        = sorted(yr_totals.index)

    df_e = df_exc.copy()
    df_e["year"]   = df_e["period"].str[:4]
    df_e["amount"] = convert_amounts(df_e["amount_jpy"].values, currency, fx) / 10_000
    yr_exc         = df_e.groupby("year")["amount"].sum()

    # KPI boxes
    box_w, box_h = Inches(2.8), Inches(1.4)
    for i, y in enumerate(years[:3]):
        x = Inches(0.3 + i * 3.2)
        sh = slide.shapes.add_textbox(x, Inches(1.0), box_w, box_h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
        tf = sh.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = f"{y}年 予測合計"
        r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = _C_BLUE
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = f"{yr_totals.get(y, 0):,.0f} {unit}"
        r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = _C_DARK
        if i > 0:
            prev = yr_totals.get(years[i - 1], 0)
            chg  = (yr_totals.get(y, 0) - prev) / prev * 100 if prev else 0
            p3 = tf.add_paragraph()
            r3 = p3.add_run(); r3.text = f"前年比 {chg:+.1f}%"
            r3.font.size = Pt(10)
            r3.font.color.rgb = _C_RED if chg > 0 else _C_GRAY

    # Impact table
    cumul = sum(yr_totals.get(y, 0) - yr_exc.get(y, 0) for y in years)
    tbl_data = [["年度", f"加味({unit})", f"非加味({unit})", f"影響額({unit})", "影響率"]]
    for y in years:
        a = yr_totals.get(y, 0); e = yr_exc.get(y, 0)
        tbl_data.append([y, f"{a:,.0f}", f"{e:,.0f}", f"{a-e:+,.0f}",
                         f"{(a-e)/e*100:+.1f}%" if e else "—"])
    tbl_data.append(["累計", "", "", f"{cumul:+,.0f}", ""])

    rows, cols = len(tbl_data), len(tbl_data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(2.6), Inches(9.4), Inches(0.4 * rows)).table
    for ri, row in enumerate(tbl_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _C_BLUE
                cell.text_frame.paragraphs[0].font.color.rgb = _C_WHITE
                cell.text_frame.paragraphs[0].font.bold = True
            elif ri == rows - 1:
                cell.text_frame.paragraphs[0].font.bold = True


def _slide_market_trend(prs, conn, ing_id, ing_name):
    from src.forecast.market_forecast import aggregate_market, MarketForecaster
    from src.forecast.penetration import PenetrationForecaster

    df_m = aggregate_market(conn, ing_id)
    mfc  = MarketForecaster(ing_id, horizon=36).fit(df_m).predict()
    fit  = PenetrationForecaster(ing_id).fit(df_m)
    pen_fit = fit.get_fit_result()
    pen_fc  = fit.predict(36)

    actual_x = list(range(len(df_m)))
    fc_x     = list(range(len(df_m), len(df_m) + 36))

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))

    # 市場規模
    ax1.plot(actual_x, df_m["total_units"].values / 1000, "o-", color=_rgb_hex(BRAND_BLUE), ms=3, label="実績")
    ax1.fill_between(fc_x, np.array(mfc.lower_units) / 1000, np.array(mfc.upper_units) / 1000,
                     alpha=0.15, color=_rgb_hex(BRAND_RED))
    ax1.plot(fc_x, np.array(mfc.forecast_units) / 1000, "--", color=_rgb_hex(BRAND_RED), label="予測")
    ax1.axvline(actual_x[-1], color=_rgb_hex(BRAND_GRAY), lw=1, ls=":")
    ax1.set_title("市場規模 推移（千本）", fontsize=11)
    ax1.set_xlabel("月次インデックス"); ax1.set_ylabel("千本"); ax1.legend(fontsize=9)
    ax1.grid(True, alpha=0.3)

    # 浸透率
    ax2.scatter(actual_x, np.array(pen_fit.actual) * 100, color=_rgb_hex(BRAND_BLUE), s=12, label="実績")
    ax2.plot(actual_x, np.array(pen_fit.fitted) * 100, "--", color=_rgb_hex(BRAND_GREEN), lw=1.2, label="フィット")
    ax2.plot(fc_x, pen_fc * 100, color=_rgb_hex(BRAND_RED), lw=2, label="予測")
    ax2.axvline(actual_x[-1], color=_rgb_hex(BRAND_GRAY), lw=1, ls=":")
    ax2.set_title(f"GE/BS 浸透率（上限 L={pen_fit.params.L:.1%}）", fontsize=11)
    ax2.set_xlabel("月次インデックス"); ax2.set_ylabel("%"); ax2.legend(fontsize=9)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax2.grid(True, alpha=0.3)

    fig.suptitle(f"疾患領域・成分トレンド  ─  {ing_name}", fontsize=12, fontweight="bold")
    plt.tight_layout()

    slide = _blank_slide(prs)
    _add_title_bar(slide, "② 疾患領域・成分トレンド", ing_name)
    img_buf = _fig_to_bytes(fig)
    slide.shapes.add_picture(img_buf, Inches(0.2), Inches(0.9), Inches(9.6), Inches(6.3))


def _slide_own_forecast(prs, df_adj, currency, fx, ing_name):
    unit = "万JPY" if currency == "JPY" else "万USD"
    df = df_adj.copy()
    df["amount"] = convert_amounts(df["amount_jpy"].values, currency, fx) / 10_000
    skus = sorted(df["sku_id"].unique())
    periods = sorted(df["period"].unique())

    colors = ["#3B82F6","#EF4444","#10B981","#F59E0B","#8B5CF6","#EC4899","#14B8A6"]
    fig, ax = plt.subplots(figsize=(11, 4.5))
    bottoms = np.zeros(len(periods))
    p_idx = {p: i for i, p in enumerate(periods)}

    for i, sku in enumerate(skus):
        sub = df[df["sku_id"] == sku].sort_values("period")
        vals = np.array([sub[sub["period"] == p]["amount"].sum() for p in periods])
        ax.bar(range(len(periods)), vals, bottom=bottoms, label=sku, color=colors[i % len(colors)])
        bottoms += vals

    tick_step = max(1, len(periods) // 8)
    ax.set_xticks(range(0, len(periods), tick_step))
    ax.set_xticklabels([periods[i] for i in range(0, len(periods), tick_step)], rotation=45, ha="right", fontsize=8)
    ax.set_ylabel(f"金額（{unit}）"); ax.set_xlabel("月")
    ax.set_title(f"自社 SKU 月次 Sell-out 予測  ─  {ing_name} ({unit})", fontsize=11)
    ax.legend(fontsize=9, loc="upper left"); ax.grid(True, axis="y", alpha=0.3)
    plt.tight_layout()

    slide = _blank_slide(prs)
    _add_title_bar(slide, "③ 自社販売予測（SKU別月次）", f"{ing_name} — {unit}")
    img_buf = _fig_to_bytes(fig)
    slide.shapes.add_picture(img_buf, Inches(0.2), Inches(0.9), Inches(9.6), Inches(6.3))


def _slide_scenario(prs, conn, ing_id, ing_name, axis2):
    from src.forecast.market_forecast import aggregate_market
    from src.forecast.penetration import PenetrationForecaster

    df_m  = aggregate_market(conn, ing_id)
    start = next_period(df_m["period"].iloc[-1])
    ev_adj = build_penetration_events(conn, start, 36, "regulatory_adjusted", axis2)
    ev_exc = build_penetration_events(conn, start, 36, "regulatory_excluded",  axis2)

    pf = PenetrationForecaster(ing_id).fit(df_m)
    pen_adj = pf.predict(36, events=ev_adj) * 100
    pen_exc = pf.predict(36, events=ev_exc) * 100

    actual_n = len(df_m)
    fc_x     = list(range(actual_n, actual_n + 36))
    actual_x = list(range(actual_n))

    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.scatter(actual_x, df_m["penetration_rate"].values * 100,
               color=_rgb_hex(BRAND_GRAY), s=12, label="実績", zorder=3)
    ax.plot(fc_x, pen_exc, color=_rgb_hex(BRAND_BLUE), lw=2, label="非加味（Organic）")
    ax.plot(fc_x, pen_adj, color=_rgb_hex(BRAND_RED),  lw=2, label="加味（Adjusted）")
    ax.fill_between(fc_x, pen_exc, pen_adj, alpha=0.15, color=_rgb_hex(BRAND_RED), label="制度変更上乗せ幅")
    ax.axvline(actual_x[-1], color=_rgb_hex(BRAND_GRAY), lw=1, ls=":")
    ax.set_xlabel("月次インデックス"); ax.set_ylabel("浸透率（%）")
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.set_title(f"制度変更影響比較（浸透率カーブ）  ─  {ing_name}", fontsize=11)
    ax.legend(fontsize=9); ax.grid(True, alpha=0.3)
    plt.tight_layout()

    slide = _blank_slide(prs)
    _add_title_bar(slide, "④ 制度変更影響比較", f"加味（Adjusted）vs 非加味（Organic） — {ing_name}")
    img_buf = _fig_to_bytes(fig)
    slide.shapes.add_picture(img_buf, Inches(0.2), Inches(0.9), Inches(9.6), Inches(6.3))


def _slide_waterfall(prs, actual_12m, fc_exc_12m, fc_adj_12m, currency, fx, ing_name):
    unit = "万JPY" if currency == "JPY" else "万USD"
    div  = 10_000

    base    = actual_12m / div
    organic = (fc_exc_12m - actual_12m) / div
    reg     = (fc_adj_12m - fc_exc_12m) / div
    total   = fc_adj_12m / div

    labels = ["前期実績（12M）", "市場成長・浸透率等\n（オーガニック）", "制度変更影響", "予測合計（12M）"]
    values = [base, organic, reg, total]
    colors_wf = [_rgb_hex(BRAND_BLUE), _rgb_hex(BRAND_GREEN) if organic >= 0 else _rgb_hex(BRAND_RED),
                 _rgb_hex(BRAND_RED) if reg >= 0 else _rgb_hex(BRAND_GREEN), _rgb_hex(BRAND_BLUE)]
    bottoms = [0, base, base + organic, 0]

    fig, ax = plt.subplots(figsize=(9, 4.5))
    for i, (lbl, val, bot, col) in enumerate(zip(labels, values, bottoms, colors_wf)):
        ax.bar(i, abs(val), bottom=bot if i < 3 else 0, color=col, width=0.5, alpha=0.85)
        sign = "+" if val >= 0 else ""
        ax.text(i, bot + (abs(val) / 2 if i < 3 else abs(val) / 2), f"{sign}{val:,.0f}",
                ha="center", va="center", fontsize=9, color="white", fontweight="bold")

    # connector lines
    running = [base, base + organic]
    for i, r in enumerate(running):
        ax.plot([i + 0.25, i + 0.75], [r, r], color=_rgb_hex(BRAND_GRAY), lw=0.8, ls="--")

    ax.set_xticks(range(4)); ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel(f"金額（{unit}）"); ax.set_title(f"要因分解（ウォーターフォール）  ─  {ing_name}", fontsize=11)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:,.0f}"))
    ax.grid(True, axis="y", alpha=0.3); plt.tight_layout()

    slide = _blank_slide(prs)
    _add_title_bar(slide, "⑤ 要因分解", f"前期実績 → 12ヶ月予測の変化要因分解  ─  {ing_name}")
    img_buf = _fig_to_bytes(fig)
    slide.shapes.add_picture(img_buf, Inches(0.5), Inches(0.95), Inches(9.0), Inches(6.1))


def _slide_risk(prs, ing_name):
    slide = _blank_slide(prs)
    _add_title_bar(slide, "⑥ リスク・アクション", ing_name)

    sections = [
        ("在庫状況",   "・卸別在庫日数: 要確認（直近データを参照）\n・薬価改定前後の積み増し・調整在庫に注意"),
        ("リスク事項", "・制度変更の施行時期が前後する可能性あり\n・競合品の新規参入状況を確認\n・為替レートの変動（USD 表示時）"),
        ("アクション", "・□ 卸別在庫データを最新版に更新\n・□ IQVIA データ最新号の取込確認\n・□ 次期薬価改定情報（実施時期・改定率）を確認"),
        ("備考",       "（手動で記入してください）"),
    ]
    y_start = 1.1
    for title, body in sections:
        tf_t = slide.shapes.add_textbox(Inches(0.3), Inches(y_start), Inches(2.0), Inches(0.4))
        tf_t.text_frame.paragraphs[0].text = title
        tf_t.text_frame.paragraphs[0].runs[0].font.bold = True
        tf_t.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        tf_t.text_frame.paragraphs[0].runs[0].font.color.rgb = _C_BLUE

        tf_b = slide.shapes.add_textbox(Inches(2.4), Inches(y_start), Inches(7.3), Inches(1.1))
        tf_b.fill.solid(); tf_b.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)
        for line in body.split("\n"):
            para = tf_b.text_frame.add_paragraph() if tf_b.text_frame.paragraphs[0].text else tf_b.text_frame.paragraphs[0]
            para.text = line
            para.font.size = Pt(10)
        y_start += 1.55


def generate_pptx(
    conn, ing_id: str, axis2: str, currency: str, fx: float
) -> bytes:
    ing_name = INGREDIENTS.get(ing_id, ing_id)

    df_adj = load_sku_fc_df(conn, ing_id, "regulatory_adjusted", axis2)
    df_exc = load_sku_fc_df(conn, ing_id, "regulatory_excluded",  axis2)

    # 前期・予測 12 ヶ月合計（ウォーターフォール用）
    actual_12m  = get_actual_sellout_12m(conn, ing_id)
    fc_adj_12m  = convert_amounts(df_adj["amount_jpy"].values[:12 * df_adj["sku_id"].nunique()], currency, fx).sum()
    fc_exc_12m  = convert_amounts(df_exc["amount_jpy"].values[:12 * df_exc["sku_id"].nunique()], currency, fx).sum()
    actual_12m_c = convert_amounts(np.array([actual_12m]), currency, fx)[0]

    prs = _new_prs()
    _slide_title(prs, ing_name, axis2, currency, str(date.today()))
    _slide_summary(prs, df_adj.copy(), df_exc.copy(), currency, fx, ing_name, axis2)
    _slide_market_trend(prs, conn, ing_id, ing_name)
    _slide_own_forecast(prs, df_adj.copy(), currency, fx, ing_name)
    _slide_scenario(prs, conn, ing_id, ing_name, axis2)
    _slide_waterfall(prs, actual_12m_c, fc_exc_12m, fc_adj_12m, currency, fx, ing_name)
    _slide_risk(prs, ing_name)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
